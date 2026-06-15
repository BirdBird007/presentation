import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

# ---------------------------------------------------------------------------
# COLOR PALETTE DEFINITIONS (Premium, Modern Theme)
# ---------------------------------------------------------------------------
COLOR_DARK_BG = RGBColor(15, 23, 42)      # Deep Slate / Dark Charcoal
COLOR_LIGHT_BG = RGBColor(248, 250, 252)  # Soft Off-White / Light Slate Gray
COLOR_CARD_BG = RGBColor(255, 255, 255)   # Pure White for panels/cards
COLOR_BORDER = RGBColor(226, 232, 240)    # Light gray border
COLOR_TEXT_PRIMARY = RGBColor(15, 23, 42) # Slate black for titles
COLOR_TEXT_MUTED = RGBColor(100, 116, 139)# Slate gray for descriptions
COLOR_GOLD = RGBColor(197, 160, 89)       # Champagne Gold (Udomkunnatum Brand)
COLOR_GOLD_DARK = RGBColor(157, 120, 49)  # Darker Gold for text contrast
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DEEP_BLUE = RGBColor(0, 51, 153)    # Corporate Deep Blue (GC Group theme)

FONT_HEADING = "Arial"
FONT_BODY = "Arial"

def create_presentation():
    # Load from the GC Group template to preserve master layouts and styles
    template_filename = "GCグループ_プレゼンテーションフォーマット.pptx"
    
    if not os.path.exists(template_filename):
        print(f"[ERROR] Template file {template_filename} not found!")
        return

    prs = Presentation(template_filename)
    
    # 1. Dynamically extract the GC Group logo picture from Slide 2 of the template
    # and save it as a local file 'gc_logo.png'
    logo_path = "gc_logo.png"
    logo_extracted = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13: # PICTURE
                try:
                    with open(logo_path, 'wb') as f:
                        f.write(shape.image.blob)
                    logo_extracted = True
                    break
                except Exception as e:
                    print(f"[WARNING] Failed to extract logo: {e}")
        if logo_extracted:
            break
            
    # 2. Delete all default slides from the loaded template to start clean
    slds = prs.slides._sldIdLst
    for i in range(len(prs.slides)-1, -1, -1):
        prs.part.drop_rel(slds[i].rId)
        slds.remove(slds[i])
        
    content_layout = prs.slide_layouts[1] # "タイトルとコンテンツ" layout (Layout 1)
    
    def add_clean_slide():
        slide = prs.slides.add_slide(content_layout)
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        return slide

    # -----------------------------------------------------------------------
    # HELPER FUNCTIONS FOR SHAPE STYLING & COMMON ELEMENTS
    # -----------------------------------------------------------------------
    def set_shape_color(shape, fill_color, border_color=None):
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()

    def add_header_and_logo(slide, title_text, subtitle_text=None):
        # Add the GC Logo in the top-left corner to match the template exactly
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(0.816), Inches(0.45), Inches(1.2), Inches(1.25))
            
        # Header text frame placed next to the logo
        title_box = slide.shapes.add_textbox(Inches(2.3), Inches(0.45), Inches(10.2), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = FONT_HEADING
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_PRIMARY
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = FONT_BODY
            p2.font.size = Pt(12)
            p2.font.color.rgb = COLOR_TEXT_MUTED
            p2.space_before = Pt(3)

    # -----------------------------------------------------------------------
    # SLIDE 1: Title Slide (White Theme matching GC Template cover)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    
    # Add GC Logo on Title slide at the template position
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.588), Inches(0.252), Inches(1.8), Inches(1.88))
        
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p0 = tf.paragraphs[0]
    p0.text = "BUSINESS PORTFOLIO & PARTNERSHIP PROPOSAL"
    p0.font.name = FONT_HEADING
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_GOLD_DARK
    p0.space_after = Pt(12)
    
    p1 = tf.add_paragraph()
    p1.text = "THE UDOMKUNNATUM FAMILY"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "ウドムクンナタム家と共同出資ビジネスグループ"
    p2.font.name = FONT_HEADING
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_DEEP_BLUE
    p2.space_before = Pt(8)
    
    p3 = tf.add_paragraph()
    p3.text = "Exploring Strategic Collaborations in Retail, Premium Real Estate, and Wellness Beverages\n日本企業との戦略的提携に向けたビジネスプロファイル"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.space_before = Pt(24)

    # -----------------------------------------------------------------------
    # SLIDE 2: Timeline of Progress (TCC "歩み" style)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Udomkunnatum Family Businessの歩み (Historical Growth)",
                        "From retail pioneers to premium real estate and wellness beverage holdings (小売・不動産からウェルネスまで)")
    
    # Left Card: Philosophy and Context
    left_x = Inches(0.75)
    left_y = Inches(1.9)
    left_w = Inches(3.6)
    left_h = Inches(4.8)
    
    card_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, left_y, left_w, left_h)
    set_shape_color(card_left, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_left = slide.shapes.add_textbox(left_x + Inches(0.3), left_y + Inches(0.3), left_w - Inches(0.6), left_h - Inches(0.6))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0
    
    p_ph_title = tf_left.paragraphs[0]
    p_ph_title.text = "FOUNDER'S PHILOSOPHY"
    p_ph_title.font.name = FONT_HEADING
    p_ph_title.font.size = Pt(13)
    p_ph_title.font.bold = True
    p_ph_title.font.color.rgb = COLOR_GOLD_DARK
    p_ph_title.space_after = Pt(10)
    
    p_val_th = tf_left.add_paragraph()
    p_val_th.text = "相互信頼と友情（クンナタム・ナムミット） & Value First"
    p_val_th.font.name = FONT_HEADING
    p_val_th.font.size = Pt(16)
    p_val_th.font.bold = True
    p_val_th.font.color.rgb = COLOR_DEEP_BLUE
    
    p_val_desc = tf_left.add_paragraph()
    p_val_desc.text = "Our core value \"クンナタム・ナムミット（相互信頼と友情）\" (Mutual Trust & Friendship) aligns closely with the Japanese spirit of \"Kyosei\" (Co-existence) and \"Sanpo Yoshi\" (Three-way benefit). We create long-term value for consumers first; sales and profit naturally follow.\n\n私たちは消費者のための価値創造を第一の目標と信じており、それが40年以上にわたる持続可能なビジネスの成長をもたらしています。"
    p_val_desc.font.name = FONT_BODY
    p_val_desc.font.size = Pt(11)
    p_val_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_val_desc.space_before = Pt(12)
    p_val_desc.line_spacing = 1.15

    # Right Timeline: Diagonal ascending line
    start_tx = Inches(4.8)
    start_ty = Inches(5.8)
    end_tx = Inches(11.5)
    end_ty = Inches(2.2)
    
    # Draw timeline diagonal line
    timeline_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_tx, start_ty, end_tx, end_ty)
    timeline_line.line.color.rgb = COLOR_DEEP_BLUE
    timeline_line.line.width = Pt(3.5)
    
    # 6 Timeline Nodes
    nodes = [
        # (x, y, year, text_en, text_th, text_pos_x, text_pos_y)
        (4.80, 5.80, "1980s", "Robinson Dept Store", "ロビンソン百貨店（共同創業者）", 4.5, 6.2),
        (6.14, 5.08, "1995", "HomePro Co-founded", "ホームプロ（LH社との共同出資）", 5.8, 3.8),
        (7.48, 4.36, "2000s", "Jungceylon Phuket", "プーケット小売商業施設管理", 7.1, 4.7),
        (8.82, 3.64, "2010s", "Swan Lake Khao Yai", "スワンレイク（4万本の植林）", 8.5, 2.4),
        (10.16, 2.92, "2023", "6ty Degrees Launch", "チェンダオ天然水", 9.8, 3.3),
        (11.50, 2.20, "Future", "Wellness Ecosystem", "300ライ規模のウェルネス総合開発", 11.2, 1.0)
    ]
    
    for i, (nx, ny, year, t_en, t_th, tx, ty) in enumerate(nodes):
        # Draw node circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - 0.22), Inches(ny - 0.22), Inches(0.44), Inches(0.44))
        set_shape_color(circle, COLOR_GOLD, COLOR_DEEP_BLUE)
        circle.line.width = Pt(2)
        
        # Draw small inner circle
        inner = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - 0.1), Inches(ny - 0.1), Inches(0.2), Inches(0.2))
        set_shape_color(inner, COLOR_DEEP_BLUE)
        
        # Connect node to text box
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(nx), Inches(ny), Inches(tx + 1.1), Inches(ty + 0.4))
        conn.line.color.rgb = COLOR_BORDER
        conn.line.width = Pt(1)
        
        # Textbox for descriptions
        tb_node = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(2.2), Inches(0.9))
        tf_node = tb_node.text_frame
        tf_node.word_wrap = True
        tf_node.margin_left = tf_node.margin_right = tf_node.margin_top = tf_node.margin_bottom = 0
        
        p_yr = tf_node.paragraphs[0]
        p_yr.text = year
        p_yr.font.name = FONT_HEADING
        p_yr.font.size = Pt(13)
        p_yr.font.bold = True
        p_yr.font.color.rgb = COLOR_GOLD_DARK
        
        p_te = tf_node.add_paragraph()
        p_te.text = t_en
        p_te.font.name = FONT_HEADING
        p_te.font.size = Pt(11)
        p_te.font.bold = True
        p_te.font.color.rgb = COLOR_TEXT_PRIMARY
        p_te.space_before = Pt(2)
        
        p_tt = tf_node.add_paragraph()
        p_tt.text = t_th
        p_tt.font.name = FONT_BODY
        p_tt.font.size = Pt(9.5)
        p_tt.font.color.rgb = COLOR_TEXT_MUTED
        p_tt.space_before = Pt(1)

    # -----------------------------------------------------------------------
    # SLIDE 3: 4-Column Business Pillars (TCC "事業" style)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Udomkunnatum Groupの事業領域 (Business Divisions)",
                        "Four core business sectors driven by ethical value and premium quality (4つのコア事業部門)")
    
    col_w = Inches(2.6)
    col_h = Inches(4.5)
    col_y = Inches(2.0)
    col_gap = Inches(0.4)
    start_x = Inches(0.75)
    
    pillars = [
        {
            "num": "01",
            "title_en": "Retail & Distribution",
            "title_th": "小売＆流通",
            "points": [
                "Home Product Center (HomePro) co-founded in 1995.",
                "Undisputed #1 home improvement network in Thailand.",
                "Over 100 branches across Thailand, Malaysia, Vietnam.",
                "Leading importer of premium home assets."
            ]
        },
        {
            "num": "02",
            "title_en": "Premium Real Estate",
            "title_th": "高級不動産",
            "points": [
                "Elysian Development focusing on luxury eco-living.",
                "Swan Lake Khao Yai: Valued at 3B+ THB.",
                "Reforested 40,000+ trees on dry land plot.",
                "80% dedicated to forest and green common areas."
            ]
        },
        {
            "num": "03",
            "title_en": "Hospitality & Lifestyle",
            "title_th": "ホスピタリティ＆ライフスタイル",
            "points": [
                "Jungceylon Patong Phuket space planning & commercial curation.",
                "Fashion apparel brand FQ&L scaled nationwide.",
                "Premium boutique hotels and lifestyle property operations."
            ]
        },
        {
            "num": "04",
            "title_en": "Wellness FMCG",
            "title_th": "ウェルネスFMCG",
            "points": [
                "Rare Beverage Co. founded for healthy products.",
                "6ty Degrees mineral water: 850k bottles/day capacity.",
                "1.0 Billion THB smart factory in Chiang Dao.",
                "Certified by SGS, Intertek, and ALS mineral labs."
            ]
        }
    ]
    
    for i, pil in enumerate(pillars):
        x = start_x + i * (col_w + col_gap)
        
        # Draw card container
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, col_y, col_w, col_h)
        set_shape_color(card, COLOR_CARD_BG, COLOR_BORDER)
        
        # Draw colored header bar inside card
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, col_y, col_w, Inches(0.4))
        set_shape_color(header_bar, COLOR_DEEP_BLUE if i % 2 == 0 else COLOR_GOLD)
        
        tb = slide.shapes.add_textbox(x + Inches(0.2), col_y + Inches(0.5), col_w - Inches(0.4), col_h - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p_num = tf.paragraphs[0]
        p_num.text = pil["num"]
        p_num.font.name = FONT_HEADING
        p_num.font.size = Pt(24)
        p_num.font.bold = True
        p_num.font.color.rgb = COLOR_GOLD_DARK if i % 2 == 0 else COLOR_DEEP_BLUE
        
        p_ten = tf.add_paragraph()
        p_ten.text = pil["title_en"]
        p_ten.font.name = FONT_HEADING
        p_ten.font.size = Pt(14)
        p_ten.font.bold = True
        p_ten.font.color.rgb = COLOR_TEXT_PRIMARY
        p_ten.space_before = Pt(4)
        
        p_tth = tf.add_paragraph()
        p_tth.text = pil["title_th"]
        p_tth.font.name = FONT_HEADING
        p_tth.font.size = Pt(11)
        p_tth.font.bold = True
        p_tth.font.color.rgb = COLOR_GOLD_DARK
        p_tth.space_before = Pt(2)
        p_tth.space_after = Pt(10)
        
        for pt in pil["points"]:
            p_pt = tf.add_paragraph()
            p_pt.text = f"• {pt}"
            p_pt.font.name = FONT_BODY
            p_pt.font.size = Pt(9.5)
            p_pt.font.color.rgb = COLOR_TEXT_MUTED
            p_pt.space_before = Pt(4)
            p_pt.line_spacing = 1.1

    # -----------------------------------------------------------------------
    # SLIDE 4: Family Succession Tree (TCC "家系譜" style)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Udomkunnatum家系譜・後継体制 (Founder & Succession)",
                        "A clear family governance transitioning to next-generation active leaders (確かな後継者体制とガバナンス)")
    
    # Left Card: Ethical & Corporate Governance Values
    left_x = Inches(0.75)
    left_y = Inches(1.9)
    left_w = Inches(3.6)
    left_h = Inches(4.8)
    
    card_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, left_y, left_w, left_h)
    set_shape_color(card_left, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_left = slide.shapes.add_textbox(left_x + Inches(0.3), left_y + Inches(0.3), left_w - Inches(0.6), left_h - Inches(0.6))
    tf_left = tb_left.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = tf_left.margin_right = tf_left.margin_top = tf_left.margin_bottom = 0
    
    p_gv = tf_left.paragraphs[0]
    p_gv.text = "FAMILY GOVERNANCE"
    p_gv.font.name = FONT_HEADING
    p_gv.font.size = Pt(13)
    p_gv.font.bold = True
    p_gv.font.color.rgb = COLOR_GOLD_DARK
    p_gv.space_after = Pt(12)
    
    p_gv_desc = tf_left.add_paragraph()
    p_gv_desc.text = "■ Stable Succession\nThe group focuses on long-term continuity (信用 - Shin'yo). Succession is pre-planned and execution-focused, avoiding fragmented leadership structures.\n\n■ Professional Grooming\nThe second generation is trained in top-tier retail spaces and scales new sectors with high agility.\n\n透明性とプロフェッショナリズムによる安定した事業承継を実現し、機動的な新規事業の開拓に注力しています。"
    p_gv_desc.font.name = FONT_BODY
    p_gv_desc.font.size = Pt(11)
    p_gv_desc.font.color.rgb = COLOR_TEXT_MUTED
    p_gv_desc.space_before = Pt(8)
    p_gv_desc.line_spacing = 1.15

    # Center: Family / Executive Tree Diagram
    # 1. Gen 1: Founder Manit Udomkunnatum
    # Profile photo
    manit_img_path = "manit_profile.png"
    if os.path.exists(manit_img_path):
        slide.shapes.add_picture(manit_img_path, Inches(4.8), Inches(1.9), Inches(1.5), Inches(2.0))
    else:
        # Draw placeholder box
        p_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(1.9), Inches(1.5), Inches(2.0))
        set_shape_color(p_box, COLOR_BORDER)
        
    card_m = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.9), Inches(4.5), Inches(2.0))
    set_shape_color(card_m, COLOR_CARD_BG, COLOR_DEEP_BLUE)
    card_m.line.width = Pt(1.5)
    
    tb_m = slide.shapes.add_textbox(Inches(6.7), Inches(2.0), Inches(4.1), Inches(1.8))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True
    tf_m.margin_left = tf_m.margin_right = tf_m.margin_top = tf_m.margin_bottom = 0
    
    pm1 = tf_m.paragraphs[0]
    pm1.text = "Manit Udomkunnatum (Founder)"
    pm1.font.name = FONT_HEADING
    pm1.font.size = Pt(15)
    pm1.font.bold = True
    pm1.font.color.rgb = COLOR_DEEP_BLUE
    
    pm2 = tf_m.add_paragraph()
    pm2.text = "マニット・ウドムクンナタム氏"
    pm2.font.name = FONT_HEADING
    pm2.font.size = Pt(12)
    pm2.font.bold = True
    pm2.font.color.rgb = COLOR_GOLD_DARK
    
    pm3 = tf_m.add_paragraph()
    pm3.text = "• Co-Founder, Robinson Department Store\n• Co-Founder & Vice Chairman, HomePro PCL\n• Leading advisor on strategic land bank development"
    pm3.font.name = FONT_BODY
    pm3.font.size = Pt(10)
    pm3.font.color.rgb = COLOR_TEXT_MUTED
    pm3.space_before = Pt(6)
    
    # 2. Gen 2: Eldest Successor Rena Udomkunnatum
    # Profile photo
    rena_img_path = "rena_profile.png"
    if os.path.exists(rena_img_path):
        slide.shapes.add_picture(rena_img_path, Inches(4.8), Inches(4.6), Inches(1.5), Inches(2.0))
    else:
        p_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.8), Inches(4.6), Inches(1.5), Inches(2.0))
        set_shape_color(p_box, COLOR_BORDER)
        
    card_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(4.6), Inches(4.5), Inches(2.0))
    set_shape_color(card_r, COLOR_CARD_BG, COLOR_GOLD)
    card_r.line.width = Pt(1.5)
    
    tb_r = slide.shapes.add_textbox(Inches(6.7), Inches(4.7), Inches(4.1), Inches(1.8))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    pr1 = tf_r.paragraphs[0]
    pr1.text = "Rena Udomkunnatum (Eldest Heir)"
    pr1.font.name = FONT_HEADING
    pr1.font.size = Pt(15)
    pr1.font.bold = True
    pr1.font.color.rgb = COLOR_GOLD_DARK
    
    pr2 = tf_r.add_paragraph()
    pr2.text = "リーナ・ウドムクンナタム氏（長女・後継者）"
    pr2.font.name = FONT_HEADING
    pr2.font.size = Pt(12)
    pr2.font.bold = True
    pr2.font.color.rgb = COLOR_DEEP_BLUE
    
    pr3 = tf_r.add_paragraph()
    pr3.text = "• CEO, Elysian Development Co., Ltd.\n• Founder & Managing Director, Rare Beverage Co., Ltd.\n• Former Curator, Jungceylon Patong Phuket"
    pr3.font.name = FONT_BODY
    pr3.font.size = Pt(10)
    pr3.font.color.rgb = COLOR_TEXT_MUTED
    pr3.space_before = Pt(6)

    # 3. Connection Arrow (Succession link)
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(5.55), Inches(3.9), Inches(5.55), Inches(4.6))
    conn.line.color.rgb = COLOR_GOLD_DARK
    conn.line.width = Pt(2.5)
    conn.line.end_arrowhead = 2 # Triangle arrowhead

    # Right Card: Professional governance comments
    right_x = Inches(11.2)
    right_y = Inches(1.9)
    right_w = Inches(1.383) # Adjust to fit 13.33 inches
    # Actually, let's make it a nice summary box on the right of the tree
    tb_summary = slide.shapes.add_textbox(Inches(11.1), Inches(1.9), Inches(1.5), Inches(4.8))
    tf_sum = tb_summary.text_frame
    tf_sum.word_wrap = True
    tf_sum.margin_left = tf_sum.margin_right = tf_sum.margin_top = tf_sum.margin_bottom = 0
    
    ps = tf_sum.paragraphs[0]
    ps.text = "GOVERNANCE"
    ps.font.name = FONT_HEADING
    ps.font.size = Pt(11)
    ps.font.bold = True
    ps.font.color.rgb = COLOR_DEEP_BLUE
    ps.space_after = Pt(8)
    
    ps_body = tf_sum.add_paragraph()
    ps_body.text = "Manit acts as visionary guide (Chairman level).\n\nRena leads operation as Chief Executive.\n\nClear roles align with Japanese corporate expectations."
    ps_body.font.name = FONT_BODY
    ps_body.font.size = Pt(9.5)
    ps_body.font.color.rgb = COLOR_TEXT_MUTED
    ps_body.line_spacing = 1.15

    # -----------------------------------------------------------------------
    # SLIDE 5: Pillar 1 - Retail & Distribution (HomePro PCL)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Business Pillar 1: Retail & Distribution (HomePro)",
                        "Unrivaled access to Thailand's top home improvements retail networks (大手小売市場を牽引)")
    
    col_w2 = Inches(5.6)
    col_h2 = Inches(4.8)
    col_y2 = Inches(1.9)
    
    # Left Card
    card_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), col_y2, col_w2, col_h2)
    set_shape_color(card_1, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_1 = slide.shapes.add_textbox(Inches(1.05), col_y2 + Inches(0.3), col_w2 - Inches(0.6), col_h2 - Inches(0.6))
    tf_1 = tb_1.text_frame
    tf_1.word_wrap = True
    tf_1.margin_left = tf_1.margin_right = tf_1.margin_top = tf_1.margin_bottom = 0
    
    p1_title = tf_1.paragraphs[0]
    p1_title.text = "Home Product Center PCL (HomePro)"
    p1_title.font.name = FONT_HEADING
    p1_title.font.size = Pt(18)
    p1_title.font.bold = True
    p1_title.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p1_th = tf_1.add_paragraph()
    p1_th.text = "ホーム・プロダクト・センター株式会社（HomePro）"
    p1_th.font.name = FONT_HEADING
    p1_th.font.size = Pt(12)
    p1_th.font.bold = True
    p1_th.font.color.rgb = COLOR_GOLD_DARK
    p1_th.space_before = Pt(2)
    p1_th.space_after = Pt(12)
    
    points_1 = [
        ("Market Leader", "Thailand's #1 home improvement retail chain. Publicly listed (SET: HMPRO) with strong financial backing."),
        ("Massive Footprint", "Over 100 branches in Thailand, plus growing footprints in Malaysia and Vietnam."),
        ("Japanese Brand Synergy", "Long-term partner distributing leading Japanese electronics, tools, and kitchenware brands to Thai consumers."),
        ("Strategic Advisory", "Mr. Manit continues to serve as Vice Chairman, steering retail strategies and partnership alignments.")
    ]
    
    for title, desc in points_1:
        p_t = tf_1.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_PRIMARY
        p_t.space_before = Pt(6)
        
        run = p_t.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED

    # Right Card: Image + retail space details
    card_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.983), col_y2, col_w2, col_h2)
    set_shape_color(card_2, COLOR_CARD_BG, COLOR_BORDER)
    
    # Side-by-side images: HomePro and Robinson legacy
    store_img = "modern_retail_store.png"
    rob_img = "robinson_retail.png"
    
    if os.path.exists(store_img):
        slide.shapes.add_picture(store_img, Inches(7.283), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
    else:
        p1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.283), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
        set_shape_color(p1, COLOR_BORDER)
        
    if os.path.exists(rob_img):
        slide.shapes.add_picture(rob_img, Inches(9.883), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
    else:
        p2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.883), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
        set_shape_color(p2, COLOR_BORDER)
    
    tb_2 = slide.shapes.add_textbox(Inches(7.283), col_y2 + Inches(2.3), Inches(5.0), Inches(2.2))
    tf_2 = tb_2.text_frame
    tf_2.word_wrap = True
    tf_2.margin_left = tf_2.margin_right = tf_2.margin_top = tf_2.margin_bottom = 0
    
    p2_t = tf_2.paragraphs[0]
    p2_t.text = "Retail Space & Commercial Curation"
    p2_t.font.name = FONT_HEADING
    p2_t.font.size = Pt(14)
    p2_t.font.bold = True
    p2_t.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2_desc = tf_2.add_paragraph()
    p2_desc.text = "Leveraging Manit's Robinson department store legacy and HomePro's massive retail footprint to scale lifestyle brands across Thai modern trade. Offers a direct pipeline for Japanese premium goods seeking rapid scale in Thailand."
    p2_desc.font.name = FONT_BODY
    p2_desc.font.size = Pt(10)
    p2_desc.font.color.rgb = COLOR_TEXT_MUTED
    p2_desc.space_before = Pt(6)
    p2_desc.line_spacing = 1.15
    
    # -----------------------------------------------------------------------
    # SLIDE 6: Pillar 2 - Eco-Luxury Real Estate (Swan Lake Khao Yai)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Business Pillar 2: Premium Real Estate (Swan Lake)",
                        "Developing world-class luxury residences integrated with nature and wellness (自然と調和したリゾート)")
    
    # Left Card
    card_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), col_y2, col_w2, col_h2)
    set_shape_color(card_1, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_1 = slide.shapes.add_textbox(Inches(1.05), col_y2 + Inches(0.3), col_w2 - Inches(0.6), col_h2 - Inches(0.6))
    tf_1 = tb_1.text_frame
    tf_1.word_wrap = True
    tf_1.margin_left = tf_1.margin_right = tf_1.margin_top = tf_1.margin_bottom = 0
    
    p1_title = tf_1.paragraphs[0]
    p1_title.text = "Elysian Development Co., Ltd."
    p1_title.font.name = FONT_HEADING
    p1_title.font.size = Pt(18)
    p1_title.font.bold = True
    p1_title.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p1_th = tf_1.add_paragraph()
    p1_th.text = "エリシアン・デベロップメント株式会社"
    p1_th.font.name = FONT_HEADING
    p1_th.font.size = Pt(12)
    p1_th.font.bold = True
    p1_th.font.color.rgb = COLOR_GOLD_DARK
    p1_th.space_before = Pt(2)
    p1_th.space_after = Pt(12)
    
    points_1 = [
        ("Niche Developer", "A boutique developer focusing strictly on high-end, premium properties with unique eco-living concepts."),
        ("Environmental Focus", "Integrating deep reforestation, organic agriculture, and clean energy solutions into premium residential projects."),
        ("Strategic Land Bank", "Holding highly valuable land bank assets in top premium destinations: Khao Yai, Chiang Mai, Bangkok, and Phuket."),
        ("Agile Management", "Led by Ms. Rena, ensuring fast decision-making and strict quality control on architectural execution.")
    ]
    
    for title, desc in points_1:
        p_t = tf_1.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_PRIMARY
        p_t.space_before = Pt(6)
        
        run = p_t.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED

    # Right Card: Image + project details
    card_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.983), col_y2, col_w2, col_h2)
    set_shape_color(card_2, COLOR_CARD_BG, COLOR_BORDER)
    
    # Side-by-side images: Swan Lake and Jungceylon Patong Phuket
    villa_img = "luxury_villa_lake.png"
    jc_img = "jungceylon_phuket.png"
    
    if os.path.exists(villa_img):
        slide.shapes.add_picture(villa_img, Inches(7.283), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
    else:
        p1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.283), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
        set_shape_color(p1, COLOR_BORDER)
        
    if os.path.exists(jc_img):
        slide.shapes.add_picture(jc_img, Inches(9.883), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
    else:
        p2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.883), col_y2 + Inches(0.3), Inches(2.4), Inches(1.8))
        set_shape_color(p2, COLOR_BORDER)
    
    tb_2 = slide.shapes.add_textbox(Inches(7.283), col_y2 + Inches(2.3), Inches(5.0), Inches(2.2))
    tf_2 = tb_2.text_frame
    tf_2.word_wrap = True
    tf_2.margin_left = tf_2.margin_right = tf_2.margin_top = tf_2.margin_bottom = 0
    
    p2_t = tf_2.paragraphs[0]
    p2_t.text = "Luxury Real Estate & Hospitality Curation"
    p2_t.font.name = FONT_HEADING
    p2_t.font.size = Pt(14)
    p2_t.font.bold = True
    p2_t.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2_desc = tf_2.add_paragraph()
    p2_desc.text = "Featuring Elysian's Swan Lake Khao Yai (3B+ THB eco-residential estate) and the massive Jungceylon Patong Phuket mall (200,000+ sqm) curated by Ms. Rena, demonstrating prime hospitality & real estate assets."
    p2_desc.font.name = FONT_BODY
    p2_desc.font.size = Pt(9.5)
    p2_desc.font.color.rgb = COLOR_TEXT_MUTED
    p2_desc.space_before = Pt(6)
    p2_desc.line_spacing = 1.15
    
    # -----------------------------------------------------------------------
    # SLIDE 7: Pillar 3 - Wellness & FMCG (Rare Beverage & 6ty Degrees)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Business Pillar 3: Wellness Beverages (6ty Degrees)",
                        "Pioneering natural mineral water brand sourced from clean geothermal spring (プレミアム天然水ブランド)")
    
    # Left Card
    card_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), col_y2, col_w2, col_h2)
    set_shape_color(card_1, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_1 = slide.shapes.add_textbox(Inches(1.05), col_y2 + Inches(0.3), col_w2 - Inches(0.6), col_h2 - Inches(0.6))
    tf_1 = tb_1.text_frame
    tf_1.word_wrap = True
    tf_1.margin_left = tf_1.margin_right = tf_1.margin_top = tf_1.margin_bottom = 0
    
    p1_title = tf_1.paragraphs[0]
    p1_title.text = "6ty Degrees Natural Mineral Water"
    p1_title.font.name = FONT_HEADING
    p1_title.font.size = Pt(18)
    p1_title.font.bold = True
    p1_title.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p1_th = tf_1.add_paragraph()
    p1_th.text = "6ty Degrees 天然ミネラルウォーター"
    p1_th.font.name = FONT_HEADING
    p1_th.font.size = Pt(12)
    p1_th.font.bold = True
    p1_th.font.color.rgb = COLOR_GOLD_DARK
    p1_th.space_before = Pt(2)
    p1_th.space_after = Pt(12)
    
    points_1 = [
        ("The Source", "Geothermal hot spring (60°C) in Chiang Dao, Chiang Mai, basalt-filtered at 300m+ depth. Certified by SGS, Intertek, and ALS."),
        ("Premium Composition", "Naturally heated to 60°C. Rich in 16 essential health-promoting minerals including silica and calcium."),
        ("Brand Position", "Premium wellness beverage brand targeting health-conscious consumers, upscale hotels, and global channels."),
        ("Export Ready", "Pure water quality meeting top international standards, prepared for worldwide shipping and premium distribution.")
    ]
    
    for title, desc in points_1:
        p_t = tf_1.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_PRIMARY
        p_t.space_before = Pt(6)
        
        run = p_t.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED

    # Right Card: Image + bottle details
    card_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.983), col_y2, col_w2, col_h2)
    set_shape_color(card_2, COLOR_CARD_BG, COLOR_BORDER)
    
    bottle_img = "mineral_water_bottle.png"
    if os.path.exists(bottle_img):
        slide.shapes.add_picture(bottle_img, Inches(7.283), col_y2 + Inches(0.3), Inches(5.0), Inches(2.2))
    else:
        img_placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.283), col_y2 + Inches(0.3), Inches(5.0), Inches(2.2))
        set_shape_color(img_placeholder, COLOR_BORDER)
        
    tb_2 = slide.shapes.add_textbox(Inches(7.283), col_y2 + Inches(2.7), Inches(5.0), Inches(1.8))
    tf_2 = tb_2.text_frame
    tf_2.word_wrap = True
    tf_2.margin_left = tf_2.margin_right = tf_2.margin_top = tf_2.margin_bottom = 0
    
    p2_t = tf_2.paragraphs[0]
    p2_t.text = "Healthy FMCG Expansion"
    p2_t.font.name = FONT_HEADING
    p2_t.font.size = Pt(14)
    p2_t.font.bold = True
    p2_t.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p2_desc = tf_2.add_paragraph()
    p2_desc.text = "Our wellness beverage expansion is anchored by the Chiang Dao basalt natural source. Sourced at 60°C, the water contains optimal minerals that guarantee an excellent, refreshing taste and therapeutic values. Certified by international labs, we are ready for global distribution."
    p2_desc.font.name = FONT_BODY
    p2_desc.font.size = Pt(10.5)
    p2_desc.font.color.rgb = COLOR_TEXT_MUTED
    p2_desc.space_before = Pt(8)
    p2_desc.line_spacing = 1.15

    # -----------------------------------------------------------------------
    # SLIDE 8: Production Facility - Chiang Dao Factory
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Chiang Dao Facility: 1 Billion Baht Smart Bottling Plant",
                        "High-speed automated clean-room bottling facility imported from Germany (最先端のボトリング工場)")
    
    # Left Card
    card_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), col_y2, col_w2, col_h2)
    set_shape_color(card_1, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_1 = slide.shapes.add_textbox(Inches(1.05), col_y2 + Inches(0.3), col_w2 - Inches(0.6), col_h2 - Inches(0.6))
    tf_1 = tb_1.text_frame
    tf_1.word_wrap = True
    tf_1.margin_left = tf_1.margin_right = tf_1.margin_top = tf_1.margin_bottom = 0
    
    p1_title = tf_1.paragraphs[0]
    p1_title.text = "Smart Bottling Plant (Rare Beverage)"
    p1_title.font.name = FONT_HEADING
    p1_title.font.size = Pt(18)
    p1_title.font.bold = True
    p1_title.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p1_th = tf_1.add_paragraph()
    p1_th.text = "6ty Degrees スマートボトリング工場"
    p1_th.font.name = FONT_HEADING
    p1_th.font.size = Pt(12)
    p1_th.font.bold = True
    p1_th.font.color.rgb = COLOR_GOLD_DARK
    p1_th.space_before = Pt(2)
    p1_th.space_after = Pt(12)
    
    points_1 = [
        ("German Technology", "High-speed clean-room bottling systems imported from Germany to guarantee zero human contamination."),
        ("Daily Output Capacity", "Currently producing 850,000 bottles per day (520ml and 1.25L bottles)."),
        ("Future Products", "Next expansion phase includes premium glass bottling lines and sparkling carbonated mineral water."),
        ("Ecosystem Site", "The 18-Rai smart factory sits within a 300-Rai wellness forest land holding, ready for holistic expansion.")
    ]
    
    for title, desc in points_1:
        p_t = tf_1.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_PRIMARY
        p_t.space_before = Pt(6)
        
        run = p_t.add_run()
        run.text = desc
        run.font.name = FONT_BODY
        run.font.size = Pt(10)
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED

    # Right Card: Certifications Table
    card_2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.983), col_y2, col_w2, col_h2)
    set_shape_color(card_2, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_2 = slide.shapes.add_textbox(Inches(7.283), col_y2 + Inches(0.3), Inches(5.0), Inches(4.2))
    tf_2 = tb_2.text_frame
    tf_2.word_wrap = True
    tf_2.margin_left = tf_2.margin_right = tf_2.margin_top = tf_2.margin_bottom = 0
    
    p2_t = tf_2.paragraphs[0]
    p2_t.text = "International Quality Certifications"
    p2_t.font.name = FONT_HEADING
    p2_t.font.size = Pt(14)
    p2_t.font.bold = True
    p2_t.font.color.rgb = COLOR_TEXT_PRIMARY
    p2_t.space_after = Pt(12)
    
    certs = [
        ("SGS (Australia)", "Verifies organic composition, mineral purity, and trace elements."),
        ("Intertek (United Kingdom)", "Certifies absolute biological safety and absence of contaminants."),
        ("ALS (Switzerland)", "Analyses therapeutic mineral content including soluble silica."),
        ("National Standards", "Full compliance with Thai FDA, GMP, and HACCP clean manufacturing.")
    ]
    
    for agency, desc in certs:
        p_cert = tf_2.add_paragraph()
        p_cert.text = f"■ {agency}"
        p_cert.font.name = FONT_HEADING
        p_cert.font.size = Pt(11.5)
        p_cert.font.bold = True
        p_cert.font.color.rgb = COLOR_DEEP_BLUE
        p_cert.space_before = Pt(8)
        
        p_cdesc = tf_2.add_paragraph()
        p_cdesc.text = desc
        p_cdesc.font.name = FONT_BODY
        p_cdesc.font.size = Pt(10)
        p_cdesc.font.color.rgb = COLOR_TEXT_MUTED
        p_cdesc.space_before = Pt(2)
        p_cdesc.line_spacing = 1.1

    # -----------------------------------------------------------------------
    # SLIDE 9: Strategic Value Proposition (2x2 Grid)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Our Strategic Value Proposition",
                        "Why the Udomkunnatum Group is the ideal local partner for Japanese business expansion")
    
    grid_w = Inches(5.6)
    grid_h = Inches(2.2)
    grid_y1 = Inches(1.9)
    grid_y2 = Inches(4.5)
    grid_x1 = Inches(0.75)
    grid_x2 = Inches(6.983)
    
    props = [
        (grid_x1, grid_y1, "01", "Unrivaled Local Distribution", "Direct gateway to Thailands top modern trade and home improvement segments via HomePro networks and deep commercial retail ties."),
        (grid_x2, grid_y1, "02", "Institutional Credibility", "40+ years of high-trust relationships with leading local banks, government regulators, and top-tier real estate groups like Land & Houses."),
        (grid_x1, grid_y2, "03", "Prime Land Assets", "Valuable land banks located in Thailands primary economic and tourist hubs (Bangkok, Khao Yai, Chiang Mai, Phuket) ready for joint development."),
        (grid_x2, grid_y2, "04", "Active, Professional Management", "Next-gen agile leadership that values Japanese business values (honoring agreements, focus on long-term stability, and high quality).")
    ]
    
    for x, y, num, title, desc in props:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, grid_w, grid_h)
        set_shape_color(card, COLOR_CARD_BG, COLOR_BORDER)
        
        tb = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), grid_w - Inches(0.6), grid_h - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p_t = tf.paragraphs[0]
        p_t.text = f"{num} | {title}"
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(15)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_GOLD_DARK
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.name = FONT_BODY
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_TEXT_MUTED
        p_d.space_before = Pt(6)
        p_d.line_spacing = 1.15

    # -----------------------------------------------------------------------
    # SLIDE 10: Collaboration Hypotheses (TCC "協業案" style)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "日本企業との協業案 (Strategic Collaboration Hypotheses)",
                        "Synergy cases combining Japanese tech and brands with our local assets and networks (具体的な協業仮説)")
    
    card_width = Inches(3.6)
    card_height = Inches(4.5)
    card_y = Inches(2.0)
    gap = Inches(0.5)
    start_x = Inches(0.75)
    
    models = [
        {
            "icon": "A",
            "title": "Real Estate & Wellness JV",
            "title_th": "不動産・ウェルネス共同事業",
            "bullets": [
                "Joint development of eco-luxury residences, retirement villages, or luxury wellness resorts in Chiang Mai / Khao Yai.",
                "Integrating Japanese construction tech, smart-home automation, and high-quality senior-care designs with our prime land bank."
            ]
        },
        {
            "icon": "B",
            "title": "FMCG Distribution & OEM",
            "title_th": "プレミアムミネラルウォーターの輸出＆OEM受託",
            "bullets": [
                "Exporting '6ty Degrees' mineral water to Japan's premium hospitality (HoReCa) and wellness markets.",
                "OEM/ODM bottling at our smart facility in Chiang Dao for Japanese healthy beverage brands seeking ASEAN scale."
            ]
        },
        {
            "icon": "C",
            "title": "Retail Channel Alliance",
            "title_th": "ホームプロを通じた日本製品の輸入・販売代理",
            "bullets": [
                "Distributing Japanese premium tools, green-energy devices, smart kitchenware, and lifestyle home goods in Thailand.",
                "Leveraging HomePro's 100+ branches as a direct, trusted launching pad with full marketing and logistics support."
            ]
        }
    ]
    
    for i, mod in enumerate(models):
        x = start_x + i * (card_width + gap)
        
        # Draw card container
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_width, card_height)
        set_shape_color(card, COLOR_CARD_BG, COLOR_BORDER)
        
        # Draw top color tag
        color_tag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_width, Inches(0.3))
        set_shape_color(color_tag, COLOR_DEEP_BLUE if i == 0 else (COLOR_GOLD if i == 1 else COLOR_TEXT_MUTED))
        
        tb = slide.shapes.add_textbox(x + Inches(0.3), card_y + Inches(0.4), card_width - Inches(0.6), card_height - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Header / Letter Box
        p_ic = tf.paragraphs[0]
        p_ic.text = f"PROPOSAL {mod['icon']}"
        p_ic.font.name = FONT_HEADING
        p_ic.font.size = Pt(12)
        p_ic.font.bold = True
        p_ic.font.color.rgb = COLOR_GOLD_DARK
        p_ic.space_after = Pt(8)
        
        # English Title
        p_title = tf.add_paragraph()
        p_title.text = mod["title"]
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_PRIMARY
        
        # Thai Title
        p_th = tf.add_paragraph()
        p_th.text = mod["title_th"]
        p_th.font.name = FONT_HEADING
        p_th.font.size = Pt(11)
        p_th.font.bold = True
        p_th.font.color.rgb = COLOR_DEEP_BLUE
        p_th.space_before = Pt(2)
        p_th.space_after = Pt(12)
        
        # Bullets
        for b in mod["bullets"]:
            p_b = tf.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.name = FONT_BODY
            p_b.font.size = Pt(9.5)
            p_b.font.color.rgb = COLOR_TEXT_MUTED
            p_b.space_before = Pt(6)
            p_b.line_spacing = 1.15

    # -----------------------------------------------------------------------
    # SLIDE 11: Target Beverage Partners & Strategic Proposals
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Target Beverage Partners & Solutions in Thailand",
                        "Target companies for production line, energy-saving material, and import-export solutions (協業・提案機会)")
    
    col_widths = [Inches(3.2), Inches(3.8), Inches(4.833)]
    col_xs = [Inches(0.75), Inches(3.95), Inches(7.75)]
    row_heights = [Inches(0.35)] + [Inches(0.80)] * 5
    
    row_ys = []
    current_y = Inches(1.75)
    for h in row_heights:
        row_ys.append(current_y)
        current_y += h
        
    headers = [
        {"en": "Target Company / Brand", "th": "対象企業・ブランド"},
        {"en": "Key Profile & Connections", "th": "企業概要＆連携ポイント"},
        {"en": "GC Proposed Solutions & Opportunities", "th": "GCが提案するソリューションと協業機会"}
    ]
    
    for col_idx, header in enumerate(headers):
        x = col_xs[col_idx]
        y = row_ys[0]
        w = col_widths[col_idx]
        h = row_heights[0]
        
        cell_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        set_shape_color(cell_bg, COLOR_DEEP_BLUE)
        
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.text = f"{header['en']} ({header['th']})"
        p.font.name = FONT_HEADING
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.LEFT
        
    rows_data = [
        [
            {
                "title_en": "Rare Beverage Co., Ltd.",
                "title_th": "レア・ベバレッジ株式会社",
                "sub_en": "6ty Degrees Premium Mineral Water",
                "sub_th": "6ty Degrees 天然ミネラルウォーター"
            },
            {
                "en": "• Founded by Udomkunnatum family. 1B THB invested in Chiang Dao forest ecosystem.\n• Advanced Germany Fully Automated technology producing 850,000 bottles/day.",
                "th": "・ウドムクンナタム家が創業。チェンダオの森林エリアに10億バーツを投資。\n・ドイツ製の完全自動化密閉システム技術を採用し、日産85万本の生産能力。"
            },
            {
                "en": "• Expansion Phase: Proposed energy/resource-saving materials/tech for upcoming Glass & Sparkling lines (hundreds of millions THB).\n• Export: GC can import 6ty Degrees to Japan.",
                "th": "・新ライン計画：今後新設予定のガラス瓶およびスパークリングラインに日本の省エネ・省資源資材・技術を提案。\n・輸出：GCが6ty Degreesを日本へ輸入し、日本国内市場開拓を支援。"
            }
        ],
        [
            {
                "title_en": "ThaiNamthip Co., Ltd.",
                "title_th": "タイナムティップ株式会社（ThaiNamthip）",
                "sub_en": "Major Beverage Bottler & Distributor",
                "sub_th": "コカ・コーラ等の大手飲料ボトラー"
            },
            {
                "en": "• Leading national beverage manufacturer.\n• Connections: Mr. Pornwut Sarasin serves as Chairman of the Board.",
                "th": "・タイ国内最大級の飲料製造・販売会社。\n・連携ポイント：ポーンウット・サラシン氏が同社取締役会長を努めている。"
            },
            {
                "en": "• Production Solutions: Propose premium energy and resource-saving materials (widely trusted in Japan) to cut long-term costs and boost efficiency.",
                "th": "・生産ライン提案：日本で実績のある高品質な省エネ・省資源資材を提案し、長期的なコスト削減と効率向上を支援。"
            }
        ],
        [
            {
                "title_en": "Thai Asia Pacific Brewery Co., Ltd.",
                "title_th": "タイ・アジア・パシフィック・ブリュワリー株式会社",
                "sub_en": "Major Brewer (Heineken & Tiger)",
                "sub_th": "大手ビールメーカー（ハイネケン、タイガー）"
            },
            {
                "en": "• Top national brewery.\n• Connections: Mr. Pornwut Sarasin serves as a Board Director.",
                "th": "・タイ国内大手のビール製造会社（ハイネケン、タイガー）。\n・連携ポイント：ポーンウット・サラシン氏が同社取締役を努めている。"
            },
            {
                "en": "• Green Solutions: Offer energy-saving equipment and eco-friendly packaging materials to enhance cost-efficiency and environmental compliance.",
                "th": "・グリーン提案：省エネ設備や環境配慮型パッケージ資材を提案し、環境基準への適合とコスト効率向上を両立。"
            }
        ],
        [
            {
                "title_en": "Birdy® Barista Brand",
                "title_th": "「Birdy® Barista」ブランド",
                "sub_en": "Premium RTD Coffee Segment",
                "sub_th": "プレミアム缶・ボトルコーヒー部門"
            },
            {
                "en": "• Expanding ready-to-drink premium coffee category under Ajinomoto.\n• Recently launched premium flavor series.",
                "th": "・味の素がプレミアム缶・ボトル（RTD）コーヒー市場攻略のために展開。\n・最近プレミアムな新フレーバーを発売し、市場での存在感を強化中。"
            },
            {
                "en": "• Quality Preservation: Offer Japanese-standard high-barrier materials and packaging solutions to preserve premium aroma and coffee taste.",
                "th": "・風味保持技術：プレミアムな香りと味わいを保持するため、日本基準のハイバリアパッケージ資材や包装ソリューションを提案。"
            }
        ],
        [
            {
                "title_en": "Chang Mineral Water Brand",
                "title_th": "「Chang Mineral Water」ブランド",
                "sub_en": "Leading Domestic Mineral Water",
                "sub_th": "象印（Chang）ブランドの人気天然水"
            },
            {
                "en": "• Prominent mineral water brand aiming for #1 position in domestic mind share.\n• Continuous international marketing campaigns.",
                "th": "・タイベブ（ThaiBev）傘下の主要ミネラルウォーターブランドで、国内トップシェアを目指す。\n・継続的なグローバルマーケティングキャンペーンを展開中。"
            },
            {
                "en": "• Scale Solutions: Offer production optimization solutions to scale output.\n• Trade Partner: Leverage GC's import-export network to co-expand global distribution.",
                "th": "・増産提案：生産能力拡大とプロセス最適化ソリューションを提案。\n・貿易連携：GCのグローバル輸出入ネットワークを活用し、海外販路の共同開発を提案。"
            }
        ]
    ]

    for row_idx, row in enumerate(rows_data, 1):
        y = row_ys[row_idx]
        h = row_heights[row_idx]
        
        bg_color = COLOR_WHITE if row_idx % 2 == 1 else RGBColor(245, 247, 250)
        
        for col_idx in range(3):
            x = col_xs[col_idx]
            w = col_widths[col_idx]
            
            cell_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
            set_shape_color(cell_box, bg_color, COLOR_BORDER)
            
            tb = slide.shapes.add_textbox(x, y, w, h)
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            
            data = row[col_idx]
            
            if col_idx == 0:
                p1 = tf.paragraphs[0]
                p1.text = data["title_en"]
                p1.font.name = FONT_HEADING
                p1.font.size = Pt(9.5)
                p1.font.bold = True
                p1.font.color.rgb = COLOR_TEXT_PRIMARY
                
                p2 = tf.add_paragraph()
                p2.text = data["title_th"]
                p2.font.name = FONT_HEADING
                p2.font.size = Pt(8.5)
                p2.font.bold = True
                p2.font.color.rgb = COLOR_DEEP_BLUE
                p2.space_before = Pt(1)
                
                p3 = tf.add_paragraph()
                p3.text = data["sub_en"]
                p3.font.name = FONT_BODY
                p3.font.size = Pt(8)
                p3.font.color.rgb = COLOR_TEXT_MUTED
                p3.space_before = Pt(2)
                
                p4 = tf.add_paragraph()
                p4.text = data["sub_th"]
                p4.font.name = FONT_BODY
                p4.font.size = Pt(7.5)
                p4.font.color.rgb = COLOR_TEXT_MUTED
                p4.space_before = Pt(1)
            else:
                lines_en = data["en"].split("\n")
                lines_th = data["th"].split("\n")
                
                first = True
                for i in range(max(len(lines_en), len(lines_th))):
                    if i < len(lines_en):
                        p_en = tf.paragraphs[0] if first else tf.add_paragraph()
                        first = False
                        p_en.text = lines_en[i]
                        p_en.font.name = FONT_BODY
                        p_en.font.size = Pt(8)
                        p_en.font.color.rgb = COLOR_TEXT_PRIMARY
                        p_en.space_before = Pt(3) if not first else Pt(0)
                        p_en.line_spacing = 1.05
                    
                    if i < len(lines_th):
                        p_th = tf.add_paragraph()
                        p_th.text = lines_th[i]
                        p_th.font.name = FONT_BODY
                        p_th.font.size = Pt(7.5)
                        p_th.font.color.rgb = COLOR_TEXT_MUTED
                        p_th.space_before = Pt(1)
                        p_th.line_spacing = 1.05

    # Bottom Banner for Joint Product Development (JPD)
    banner_y = Inches(6.25)
    banner_h = Inches(0.60)
    banner_w = Inches(11.833)
    banner_x = Inches(0.75)
    
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, banner_x, banner_y, banner_w, banner_h)
    set_shape_color(banner, COLOR_CARD_BG, COLOR_GOLD)
    
    tb_b = slide.shapes.add_textbox(banner_x, banner_y, banner_w, banner_h)
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = Inches(0.12)
    tf_b.margin_right = Inches(0.12)
    tf_b.margin_top = Inches(0.04)
    tf_b.margin_bottom = Inches(0.04)
    
    p_b1 = tf_b.paragraphs[0]
    p_b1.text = "Strategic Recommendation (Joint Product Development): "
    p_b1.font.name = FONT_HEADING
    p_b1.font.size = Pt(9.5)
    p_b1.font.bold = True
    p_b1.font.color.rgb = COLOR_GOLD_DARK
    
    run_en = p_b1.add_run()
    run_en.text = "Partnering to import premium raw materials from Japan for Joint Product Development (JPD) is a key entry strategy."
    run_en.font.name = FONT_BODY
    run_en.font.size = Pt(9.0)
    run_en.font.bold = False
    run_en.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p_b2 = tf_b.add_paragraph()
    p_b2.text = "追加推奨事項（共同製品開発）：タイ国内メーカーとの新飲料製品の共同開発（JPD）に向け、日本からの高品質な原材料輸入のパートナーシップ提案が効果的な参入戦略。"
    p_b2.font.name = FONT_BODY
    p_b2.font.size = Pt(8.5)
    p_b2.font.color.rgb = COLOR_TEXT_MUTED
    p_b2.space_before = Pt(1)

    # -----------------------------------------------------------------------
    # SLIDE 12: Partnership Opportunities with Aitech (A.I. Technology)
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Strategic Partnership: Aitech (A.I. Technology) (戦略的提携)",
                        "Opportunities in green manufacturing and Japanese industrial sourcing (グリーン製造・調達での協業)")
    
    # Left Card - Aitech Profile
    card_1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), col_y2, col_w2, col_h2)
    set_shape_color(card_1, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_1 = slide.shapes.add_textbox(Inches(1.05), col_y2 + Inches(0.3), col_w2 - Inches(0.6), col_h2 - Inches(0.6))
    tf_1 = tb_1.text_frame
    tf_1.word_wrap = True
    tf_1.margin_left = tf_1.margin_right = tf_1.margin_top = tf_1.margin_bottom = 0
    
    p1_title = tf_1.paragraphs[0]
    p1_title.text = "Aitech Company Profile"
    p1_title.font.name = FONT_HEADING
    p1_title.font.size = Pt(18)
    p1_title.font.bold = True
    p1_title.font.color.rgb = COLOR_TEXT_PRIMARY
    
    p1_jp = tf_1.add_paragraph()
    p1_jp.text = "Aitech 企業概要・特徴"
    p1_jp.font.name = FONT_HEADING
    p1_jp.font.size = Pt(12)
    p1_jp.font.bold = True
    p1_jp.font.color.rgb = COLOR_GOLD_DARK
    p1_jp.space_before = Pt(2)
    p1_jp.space_after = Pt(12)
    
    points_1 = [
        ("Automation Pioneer", "Leading Thai automation designer and builder for over 30 years.", 
         "30年以上の実績を持つタイの代表的な自動化・ロボットシステムインテグレーター（SI）。"),
        ("SCG Group Capital", "SCG Group (Cementhai Holding) acquired a 51% stake in 2021 to drive Smart Factory 4.0 expansion.", 
         "2021年にSCGグループ（セメントハイ・ホールディングス）が株式の51%を取得し、スマートファクトリー4.0事業を強化。"),
        ("ESG & Net Zero", "Led by President Kulchoke Phiphatana-chaichan, with a strong commitment to Green Manufacturing, ESG, and Net Zero emissions.", 
         "クン・クンチョーク社長は、グリーン製造、ESG、および産業部門のネットゼロ排出（炭素削減）を極めて重視。")
    ]
    
    for title, desc_en, desc_jp in points_1:
        p_t = tf_1.add_paragraph()
        p_t.text = f"• {title}: "
        p_t.font.name = FONT_HEADING
        p_t.font.size = Pt(10.5)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_TEXT_PRIMARY
        p_t.space_before = Pt(5)
        
        run = p_t.add_run()
        run.text = desc_en
        run.font.name = FONT_BODY
        run.font.size = Pt(9.5)
        run.font.bold = False
        run.font.color.rgb = COLOR_TEXT_MUTED
        
        p_jp = tf_1.add_paragraph()
        p_jp.text = desc_jp
        p_jp.font.name = FONT_BODY
        p_jp.font.size = Pt(8.5)
        p_jp.font.color.rgb = COLOR_TEXT_MUTED
        p_jp.space_before = Pt(1)
        p_jp.space_after = Pt(4)

    # Right Stacked Cards
    card_r1_y = col_y2
    card_r1_h = Inches(2.25)
    card_r2_y = card_r1_y + card_r1_h + Inches(0.3)
    card_r2_h = Inches(2.25)
    
    # Card R1: Green Manufacturing Partnership
    card_r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.983), card_r1_y, col_w2, card_r1_h)
    set_shape_color(card_r1, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_r1 = slide.shapes.add_textbox(Inches(7.283), card_r1_y + Inches(0.2), col_w2 - Inches(0.5), card_r1_h - Inches(0.4))
    tf_r1 = tb_r1.text_frame
    tf_r1.word_wrap = True
    tf_r1.margin_left = tf_r1.margin_right = tf_r1.margin_top = tf_r1.margin_bottom = 0
    
    p_r1_title = tf_r1.paragraphs[0]
    p_r1_title.text = "Pillar 1: Green Manufacturing Partnership"
    p_r1_title.font.name = FONT_HEADING
    p_r1_title.font.size = Pt(14)
    p_r1_title.font.bold = True
    p_r1_title.font.color.rgb = COLOR_DEEP_BLUE
    
    p_r1_jp_title = tf_r1.add_paragraph()
    p_r1_jp_title.text = "協業の柱1：グリーン製造共同提案"
    p_r1_jp_title.font.name = FONT_HEADING
    p_r1_jp_title.font.size = Pt(10)
    p_r1_jp_title.font.bold = True
    p_r1_jp_title.font.color.rgb = COLOR_GOLD_DARK
    p_r1_jp_title.space_before = Pt(1)
    p_r1_jp_title.space_after = Pt(6)
    
    p_r1_en = tf_r1.add_paragraph()
    p_r1_en.text = "Integrate GC's resource-saving and energy-saving materials into Aitech's turnkey production lines for F&B and automotive sectors to enhance eco-friendly sales points."
    p_r1_en.font.name = FONT_BODY
    p_r1_en.font.size = Pt(9.5)
    p_r1_en.font.color.rgb = COLOR_TEXT_PRIMARY
    p_r1_en.line_spacing = 1.15
    
    p_r1_jp = tf_r1.add_paragraph()
    p_r1_jp.text = "Aitechが提供する自動車および食品・飲料（F&B）分野のターンキー生産ラインに、GCの「省エネルギー・省資源資材」を組み込み、環境配慮型ソリューションとしての強みを訴求。"
    p_r1_jp.font.name = FONT_BODY
    p_r1_jp.font.size = Pt(8.5)
    p_r1_jp.font.color.rgb = COLOR_TEXT_MUTED
    p_r1_jp.space_before = Pt(4)
    p_r1_jp.line_spacing = 1.1

    # Card R2: Japanese Sourcing & Trading Partner
    card_r2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.983), card_r2_y, col_w2, card_r2_h)
    set_shape_color(card_r2, COLOR_CARD_BG, COLOR_BORDER)
    
    tb_r2 = slide.shapes.add_textbox(Inches(7.283), card_r2_y + Inches(0.2), col_w2 - Inches(0.5), card_r2_h - Inches(0.4))
    tf_r2 = tb_r2.text_frame
    tf_r2.word_wrap = True
    tf_r2.margin_left = tf_r2.margin_right = tf_r2.margin_top = tf_r2.margin_bottom = 0
    
    p_r2_title = tf_r2.paragraphs[0]
    p_r2_title.text = "Pillar 2: Japanese Sourcing & Trading Partner"
    p_r2_title.font.name = FONT_HEADING
    p_r2_title.font.size = Pt(14)
    p_r2_title.font.bold = True
    p_r2_title.font.color.rgb = COLOR_DEEP_BLUE
    
    p_r2_jp_title = tf_r2.add_paragraph()
    p_r2_jp_title.text = "協業の柱2：日本からの部材調達・商社提携"
    p_r2_jp_title.font.name = FONT_HEADING
    p_r2_jp_title.font.size = Pt(10)
    p_r2_jp_title.font.bold = True
    p_r2_jp_title.font.color.rgb = COLOR_GOLD_DARK
    p_r2_jp_title.space_before = Pt(1)
    p_r2_jp_title.space_after = Pt(6)
    
    p_r2_en = tf_r2.add_paragraph()
    p_r2_en.text = "Leverage GC's business network in Japan to source machinery parts, components, and new technologies (e.g., for brands Aitech distributes like YAMADA, KANTO SEIKI) at competitive costs."
    p_r2_en.font.name = FONT_BODY
    p_r2_en.font.size = Pt(9.5)
    p_r2_en.font.color.rgb = COLOR_TEXT_PRIMARY
    p_r2_en.line_spacing = 1.15
    
    p_r2_jp = tf_r2.add_paragraph()
    p_r2_jp.text = "Aitechの商社部門が代理店を務める日本ブランド（大智電通、ヤマダ、広高、関東精機など）の部品や、新しい自動化技術の調達において、GCの日本ネットワークを活用したコスト競争力のある迅速な供給体制を提案。"
    p_r2_jp.font.name = FONT_BODY
    p_r2_jp.font.size = Pt(8.5)
    p_r2_jp.font.color.rgb = COLOR_TEXT_MUTED
    p_r2_jp.space_before = Pt(4)
    p_r2_jp.line_spacing = 1.1

    # -----------------------------------------------------------------------
    # SLIDE 13: Other Industrial Automation & Robot SI Partners
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "Other Automation & Robot SI Partners in Thailand (その他SIパートナー)",
                        "Strategic collaboration models with specialized Japanese and local SIs (専門分野別協業モデル)")
    
    card_width = Inches(3.6)
    card_height = Inches(4.8)
    card_y = Inches(1.9)
    gap = Inches(0.5)
    start_x = Inches(0.75)
    
    si_partners = [
        {
            "icon": "1",
            "title": "KUROITSU (Thailand)",
            "sub": "Japanese Sourcing Alliance",
            "sub_jp": "日系専用機メーカー調達提携",
            "desc_en": "Japanese automation and special machine manufacturer with 13+ years in Thailand. Focuses 90% on automotive clients. GC can supply Japanese-standard parts and high-precision materials.",
            "desc_jp": "タイで13年の実績を持つ日系SI・専用機メーカー。顧客の90%が自動車業界。高い精度（High Precision）が求められる機械組立に向け、日本基準を満たすGCの高品質部材や部品供給の連携を提案。"
        },
        {
            "icon": "2",
            "title": "AEP & DDD Auto",
            "sub": "F&B Logistics Solutions",
            "sub_jp": "食品・包装・物流トータル提案",
            "desc_en": "Asia Engineering Pac (AEP) (40+ years in F&B packaging) and DDD Auto (AGV/AMR, ASRS). Partner to offer total solutions combining GC's energy-saving materials with their conveyors/robotics.",
            "desc_jp": "40年の実績を持つ食品・包装機械大手AEPと、無人搬送車（AGV/AMR）や自動倉庫（ASRS）に強みを持つDDD Auto。GCの省エネ資材と両社のコンベア・搬送システムを組み合わせ、F&B業界向けトータルソリューションを共同提案。"
        },
        {
            "icon": "3",
            "title": "Eureka & Mygrowtech",
            "sub": "Smart Warehousing & Energy Saving",
            "sub_jp": "スマート倉庫の省エネルギー支援",
            "desc_en": "Specialists in Robot Palletizer, Smart Logistics, and Vision inspection. Integrate GC's energy-saving solutions into automated warehouses and smart conveyors to reduce long-term costs for public enterprises.",
            "desc_jp": "ロボットパレタイザー、スマート物流、画像検査システムに強み。大企業や上場企業の工場・倉庫向け自動化ラインにGCの省エネソリューションを導入し、顧客の長期的な光熱費・操業コスト削減を共同で支援。"
        }
    ]
    
    for i, si in enumerate(si_partners):
        x = start_x + i * (card_width + gap)
        
        # Draw card container
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_width, card_height)
        set_shape_color(card, COLOR_CARD_BG, COLOR_BORDER)
        
        # Draw top color tag
        color_tag = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, card_y, card_width, Inches(0.3))
        set_shape_color(color_tag, COLOR_DEEP_BLUE if i == 0 else (COLOR_GOLD if i == 1 else COLOR_TEXT_MUTED))
        
        tb = slide.shapes.add_textbox(x + Inches(0.25), card_y + Inches(0.4), card_width - Inches(0.5), card_height - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Header / Letter Box
        p_ic = tf.paragraphs[0]
        p_ic.text = f"PARTNER {si['icon']}"
        p_ic.font.name = FONT_HEADING
        p_ic.font.size = Pt(11)
        p_ic.font.bold = True
        p_ic.font.color.rgb = COLOR_GOLD_DARK
        p_ic.space_after = Pt(4)
        
        # English Title
        p_title = tf.add_paragraph()
        p_title.text = si["title"]
        p_title.font.name = FONT_HEADING
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_PRIMARY
        
        # Subtitle EN
        p_sub = tf.add_paragraph()
        p_sub.text = si["sub"]
        p_sub.font.name = FONT_HEADING
        p_sub.font.size = Pt(9.5)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_DEEP_BLUE
        p_sub.space_before = Pt(2)
        
        # Subtitle JP
        p_sub_jp = tf.add_paragraph()
        p_sub_jp.text = si["sub_jp"]
        p_sub_jp.font.name = FONT_HEADING
        p_sub_jp.font.size = Pt(8.5)
        p_sub_jp.font.bold = True
        p_sub_jp.font.color.rgb = COLOR_GOLD_DARK
        p_sub_jp.space_before = Pt(1)
        p_sub_jp.space_after = Pt(10)
        
        # English Description
        p_desc = tf.add_paragraph()
        p_desc.text = si["desc_en"]
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(9.0)
        p_desc.font.color.rgb = COLOR_TEXT_PRIMARY
        p_desc.line_spacing = 1.15
        
        # Japanese Description
        p_desc_jp = tf.add_paragraph()
        p_desc_jp.text = si["desc_jp"]
        p_desc_jp.font.name = FONT_BODY
        p_desc_jp.font.size = Pt(8.0)
        p_desc_jp.font.color.rgb = COLOR_TEXT_MUTED
        p_desc_jp.space_before = Pt(6)
        p_desc_jp.line_spacing = 1.1

    # -----------------------------------------------------------------------
    # SLIDE 14: References & Sources
    # -----------------------------------------------------------------------
    slide = add_clean_slide()
    add_header_and_logo(slide, "References & Data Sources (参考文献・資料出典)",
                        "Citations and primary verification sources for all facts and figures (データの根拠)")
    
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.7), Inches(11.833), Inches(5.1))
    set_shape_color(card, COLOR_CARD_BG, COLOR_BORDER)
    
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.9), Inches(11.3), Inches(4.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_ref_title = tf.paragraphs[0]
    p_ref_title.text = "PRIMARY DOCUMENTATION & IMAGE REFERENCE SOURCES:"
    p_ref_title.font.name = FONT_HEADING
    p_ref_title.font.size = Pt(14)
    p_ref_title.font.bold = True
    p_ref_title.font.color.rgb = COLOR_GOLD_DARK
    p_ref_title.space_after = Pt(10)
    
    p_ref_1 = tf.add_paragraph()
    p_ref_1.text = "1. \"ウドムクンナタム家の後継者、ファミリービジネスの帝国をミネラルウォーターへと拡大\" [Udomkunnatum Heir Expands Family Empire], Bangkok Biz News. Published August 23, 2023."
    p_ref_1.font.name = FONT_HEADING
    p_ref_1.font.size = Pt(11)
    p_ref_1.font.bold = True
    p_ref_1.font.color.rgb = COLOR_TEXT_PRIMARY
    p_ref_1.space_before = Pt(6)
    
    p_ref_1_sub = tf.add_paragraph()
    p_ref_1_sub.text = (
        "Source Link: https://www.bangkokbiznews.com/business/business/1084892 (Accessed June 8, 2026).\n"
        "Verifies the 1.0 Billion Baht investment, 300-Rai wellness land ecosystem in Chiang Dao, German Krones tech, and Swan Lake project facts.\n"
        "Sourced Real Images:\n"
        "  • Mr. Manit Profile: https://image.bangkokbiznews.com/uploads/images/contents/w1024/2023/08/0VKYtVYJMd7B6BykqtTD.webp\n"
        "  • Ms. Rena Profile: https://image.bangkokbiznews.com/uploads/images/contents/w1024/2023/08/J2LP0qT4Q7Kkg2PbQJhM.webp\n"
        "  • Swan Lake Khao Yai: https://image.bangkokbiznews.com/uploads/images/contents/w1024/2023/08/22NHoKbGdtN3hgCtCLCf.webp\n"
        "  • 6ty Degrees Bottles: https://image.bangkokbiznews.com/uploads/images/md/2023/08/kwmcRGujNUKZM65EEgRT.webp"
    )
    p_ref_1_sub.font.name = FONT_BODY
    p_ref_1_sub.font.size = Pt(9.5)
    p_ref_1_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_ref_1_sub.space_before = Pt(2)
    p_ref_1_sub.line_spacing = 1.1
    
    p_ref_2 = tf.add_paragraph()
    p_ref_2.text = "2. Home Product Center Public Company Limited (HomePro) Annual Financial Statements & SET Disclosure (SET: HMPRO)."
    p_ref_2.font.name = FONT_HEADING
    p_ref_2.font.size = Pt(11)
    p_ref_2.font.bold = True
    p_ref_2.font.color.rgb = COLOR_TEXT_PRIMARY
    p_ref_2.space_before = Pt(8)
    
    p_ref_2_sub = tf.add_paragraph()
    p_ref_2_sub.text = "Verifies retail distribution capacity, Vice Chairman role of Mr. Manit Udomkunnatum, and partnerships with leading Japanese hardware brands."
    p_ref_2_sub.font.name = FONT_BODY
    p_ref_2_sub.font.size = Pt(9.5)
    p_ref_2_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_ref_2_sub.space_before = Pt(2)
    
    p_ref_3 = tf.add_paragraph()
    p_ref_3.text = "3. \"ホームプロが新サービス「SHOP4YOU」を開始\" [HomePro Launches New SHOP4YOU Service], Brand Inside. Published March 26, 2020."
    p_ref_3.font.name = FONT_HEADING
    p_ref_3.font.size = Pt(11)
    p_ref_3.font.bold = True
    p_ref_3.font.color.rgb = COLOR_TEXT_PRIMARY
    p_ref_3.space_before = Pt(8)
    
    p_ref_3_sub = tf.add_paragraph()
    p_ref_3_sub.text = (
        "Source Link: https://brandinside.asia/homepro-launch-new-shop4you-service/ (Accessed June 8, 2026).\n"
        "Verifies digital retail pipelines.\n"
        "Sourced Real Image:\n"
        "  • HomePro Storefront Branch: https://assets.brandinside.asia/uploads/2020/03/banner-home-01.jpg"
    )
    p_ref_3_sub.font.name = FONT_BODY
    p_ref_3_sub.font.size = Pt(9.5)
    p_ref_3_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_ref_3_sub.space_before = Pt(2)
    p_ref_3_sub.line_spacing = 1.1

    p_ref_4 = tf.add_paragraph()
    p_ref_4.text = "4. Wikimedia Commons - Robinson Department Store, Bangrak Branch storefront image."
    p_ref_4.font.name = FONT_HEADING
    p_ref_4.font.size = Pt(11)
    p_ref_4.font.bold = True
    p_ref_4.font.color.rgb = COLOR_TEXT_PRIMARY
    p_ref_4.space_before = Pt(8)
    
    p_ref_4_sub = tf.add_paragraph()
    p_ref_4_sub.text = (
        "Source Link: https://upload.wikimedia.org/wikipedia/commons/0/06/Robinson_Bangrak2019.jpg (Accessed June 8, 2026).\n"
        "Sourced Real Image:\n"
        "  • Robinson Bangrak Storefront: robinson_retail.png"
    )
    p_ref_4_sub.font.name = FONT_BODY
    p_ref_4_sub.font.size = Pt(9.5)
    p_ref_4_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_ref_4_sub.space_before = Pt(2)
    p_ref_4_sub.line_spacing = 1.1

    p_ref_5 = tf.add_paragraph()
    p_ref_5.text = "5. Wikimedia Commons - Jungceylon Shopping Center storefront image."
    p_ref_5.font.name = FONT_HEADING
    p_ref_5.font.size = Pt(11)
    p_ref_5.font.bold = True
    p_ref_5.font.color.rgb = COLOR_TEXT_PRIMARY
    p_ref_5.space_before = Pt(8)
    
    p_ref_5_sub = tf.add_paragraph()
    p_ref_5_sub.text = (
        "Source Link: https://upload.wikimedia.org/wikipedia/commons/2/2e/Jungceylon_Shopping_Center.jpg (Accessed June 8, 2026).\n"
        "Sourced Real Image:\n"
        "  • Jungceylon Patong Mall: jungceylon_phuket.png"
    )
    p_ref_5_sub.font.name = FONT_BODY
    p_ref_5_sub.font.size = Pt(9.5)
    p_ref_5_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_ref_5_sub.space_before = Pt(2)
    p_ref_5_sub.line_spacing = 1.1

    # -----------------------------------------------------------------------
    # POST-PROCESS: ENFORCE ARIAL & SCALE UP ALL FONT SIZES BY +1.5 PT
    # -----------------------------------------------------------------------
    print("[POST-PROCESS] Applying Arial font and scaling up text sizes by +1.5pt...")
    
    def adjust_font_properties(p, p_size):
        p.font.name = "Arial"
        if p_size is not None:
            p.font.size = Pt(p_size + 1.5)
        for run in p.runs:
            run.font.name = "Arial"
            run_size = run.font.size.pt if run.font.size else None
            if run_size is not None:
                run.font.size = Pt(run_size + 1.5)
            elif p_size is not None:
                run.font.size = Pt(p_size + 1.5)

    def process_shape_font(shape):
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p_size = p.font.size.pt if p.font.size else None
                adjust_font_properties(p, p_size)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for p in cell.text_frame.paragraphs:
                        p_size = p.font.size.pt if p.font.size else None
                        adjust_font_properties(p, p_size)
        elif shape.shape_type == 6: # Group shape (integer value 6)
            for sub_shape in shape.shapes:
                process_shape_font(sub_shape)

    for slide in prs.slides:
        for shape in slide.shapes:
            process_shape_font(shape)

    # -----------------------------------------------------------------------
    # SAVE PRESENTATION
    # -----------------------------------------------------------------------
    output_filename = "Udomkunnatum_Family_Business_Profile.pptx"
    prs.save(output_filename)
    print(f"[SUCCESS] PowerPoint presentation generated: {output_filename}")

if __name__ == "__main__":
    create_presentation()
