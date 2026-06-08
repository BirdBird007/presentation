import os
from pptx import Presentation

def rgb_to_hex(color):
    if color and hasattr(color, 'rgb') and color.rgb:
        return f"RGB({color.rgb[0]}, {color.rgb[1]}, {color.rgb[2]})"
    return "Unknown/None"

def inspect_file(filename, output_file):
    output_file.write(f"\n====================================================================\n")
    output_file.write(f"INSPECTING FILE: {filename}\n")
    output_file.write(f"====================================================================\n")
    
    if not os.path.exists(filename):
        output_file.write("File not found!\n")
        return
        
    prs = Presentation(filename)
    output_file.write(f"Total slides: {len(prs.slides)}\n\n")
    
    # Check Layouts
    output_file.write("--- Slide Layouts available in this template: ---\n")
    for i, layout in enumerate(prs.slide_layouts):
        output_file.write(f"Layout {i}: Name: {layout.name}\n")
        
    output_file.write("\n--- Slide Content and Shape Details: ---\n")
    for slide_idx, slide in enumerate(prs.slides):
        output_file.write(f"\nSlide {slide_idx + 1}:\n")
        output_file.write(f"  Layout name: {slide.slide_layout.name}\n")
        output_file.write(f"  Shapes found: {len(slide.shapes)}\n")
        
        for shape_idx, shape in enumerate(slide.shapes):
            output_file.write(f"  - Shape {shape_idx + 1}: Name: {shape.name}, Type: {shape.shape_type}\n")
            if shape.has_text_frame:
                tf = shape.text_frame
                for p_idx, p in enumerate(tf.paragraphs):
                    p_text = p.text.strip()
                    if p_text:
                        output_file.write(f"      Paragraph {p_idx+1}: \"{p_text}\"\n")
                        if p.font:
                            output_file.write(f"        Font: {p.font.name}, Size: {p.font.size}, Color: {rgb_to_hex(p.font.color)}\n")
            
            # Check shape fill color if applicable
            if hasattr(shape, 'fill') and shape.fill.type == 1: # solid fill
                output_file.write(f"      Fill Color: {rgb_to_hex(shape.fill.fore_color)}\n")

def main():
    with open('pptx_inspection.txt', 'w', encoding='utf-8') as f_out:
        inspect_file("GCグループ_プレゼンテーションフォーマット.pptx", f_out)
        inspect_file("TCCグループ資料　(提出版).pptx", f_out)
    print("[SUCCESS] PPTX inspection completed.")

if __name__ == "__main__":
    main()
